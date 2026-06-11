#!/usr/bin/env python3
"""Generate a professional PowerPoint presentation for Baligh project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────────────
GREEN = RGBColor(0x1B, 0x6B, 0x2F)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_DARK = RGBColor(0x0D, 0x47, 0x0A)
YELLOW = RGBColor(0xFD, 0xD8, 0x35)
YELLOW_LIGHT = RGBColor(0xFF, 0xF8, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MED = RGBColor(0xDD, 0xDD, 0xDD)
BLUE = RGBColor(0x15, 0x65, 0xC0)
BLUE_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_LIGHT = RGBColor(0xF3, 0xE5, 0xF5)

CAPTURES = "captures-application"
DIAGRAMS = "diagrammes"
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# Remove all default layouts - we use blank
blank_layout = prs.slide_layouts[6]

slide_num = [0]


def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    slide_num[0] += 1
    return slide


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect_sharp(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, font_size=14, color=DARK, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    tf.paragraphs[0].space_before = Pt(0)
    return txBox


def add_header_bar(slide, title, subtitle=None):
    add_rect_sharp(slide, Inches(0), Inches(0), W, Inches(1.3), fill_color=GREEN)
    add_text(slide, title, Inches(0.6), Inches(0.2), Inches(12), Inches(0.9),
             font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
                 font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))


def add_slide_number(slide):
    add_text(slide, str(slide_num[0]), Inches(12.4), Inches(7.0), Inches(0.7), Inches(0.4),
             font_size=9, color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_badge(slide, text, left, top, width, height, color=GREEN, font_color=WHITE, font_size=11):
    shape = add_rect(slide, left, top, width, height, fill_color=color)
    shape.adjustments[0] = 0.3
    add_text(slide, text, left, top, width, height, font_size=font_size,
             color=font_color, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_img(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {'left': left, 'top': top}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        slide.shapes.add_picture(path, **kwargs)


def add_card(slide, left, top, width, height, fill=GRAY_LIGHT):
    shape = add_rect(slide, left, top, width, height, fill_color=fill)
    shape.shadow.inherit = False
    return shape


def add_multiline(slide, lines, left, top, width, height, font_size=12, color=DARK, bold=False,
                  alignment=PP_ALIGN.LEFT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = alignment
        p.space_after = Pt(2)
    return txBox


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, GREEN)
# Decorative top stripe
add_rect_sharp(s, Inches(0), Inches(0), W, Inches(0.08), fill_color=YELLOW)
# Title area
add_text(s, "Baligh", Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.5),
         font_size=66, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "Application Mobile de Signalement Citoyen", Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.8),
         font_size=28, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text(s, "Flutter · Supabase · OpenStreetMap", Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.6),
         font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
# Separator line
add_rect_sharp(s, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.05), fill_color=YELLOW)
add_text(s, "Institut Supérieur du Numérique (SUPNUM)", Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.5),
         font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(s, "Projet de Développement Mobile — L2 DWM", Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.5),
         font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)
# Bottom
add_text(s, "Année universitaire 2025–2026", Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
         font_size=12, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)
add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Sommaire")
items = [
    ("1", "Équipe de développement"),
    ("2", "Contexte et problématique"),
    ("3", "Objectifs du projet"),
    ("4", "Architecture MVC"),
    ("5", "Stack technologique"),
    ("6", "Base de données"),
    ("7", "Fonctionnalités développées"),
    ("8", "Captures d'écran"),
    ("9", "Démonstration"),
]
for i, (num, title) in enumerate(items):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.2)
    y = Inches(1.8 + row * 1.7)
    add_card(s, x, y, Inches(3.8), Inches(1.3))
    add_circle(s, Inches(x.inches + 0.25), Inches(y.inches + 0.2), Inches(0.7), fill_color=GREEN)
    add_text(s, num, Inches(x.inches + 0.25), Inches(y.inches + 0.25), Inches(0.7), Inches(0.6),
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, Inches(x.inches + 1.15), Inches(y.inches + 0.2), Inches(2.4), Inches(0.9),
             font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — ÉQUIPE
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Équipe de Développement", "4 étudiants · L2 Développement Web et Mobile · SUPNUM")

team = [
    ("Abdy Mohameden", "24068", "Backend Supabase\nDAOs & Services", "assets/app_icon.png"),
    ("Abdselam\nAbdelvetah", "24139", "UI/UX Design\nCartographie & i18n", "assets/app_icon.png"),
    ("Hassen Oumeiry", "24238", "Gestion d'état\nContrôleurs Provider", "assets/app_icon.png"),
    ("Ahmedou Bedou", "24157", "Votes, Messagerie\nNotifications & Chat", "assets/app_icon.png"),
]
for i, (name, mat, role, icon) in enumerate(team):
    x = Inches(0.5 + i * 3.15)
    y = Inches(1.7)
    add_card(s, x, y, Inches(2.85), Inches(4.0))
    # Avatar circle
    icon_path = os.path.join(CAPTURES, supum if 'supum' in str(icon) else icon) if icon else None
    if icon and os.path.exists(icon):
        add_img(s, icon, Inches(x.inches + 0.7), Inches(y.inches + 0.25), width=Inches(1.45), height=Inches(1.45))
    else:
        add_circle(s, Inches(x.inches + 0.7), Inches(y.inches + 0.25), Inches(1.45), fill_color=GREEN)
        add_text(s, name[0], Inches(x.inches + 0.7), Inches(y.inches + 0.25), Inches(1.45), Inches(1.45),
                 font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, name, Inches(x.inches + 0.15), Inches(y.inches + 1.9), Inches(2.55), Inches(0.8),
             font_size=15, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, f"Matricule: {mat}", Inches(x.inches + 0.15), Inches(y.inches + 2.5), Inches(2.55), Inches(0.4),
             font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text(s, role, Inches(x.inches + 0.15), Inches(y.inches + 2.9), Inches(2.55), Inches(0.8),
             font_size=12, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — CONTEXTE
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Contexte & Problématique")
# Left column
add_card(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.0), fill=GREEN_LIGHT)
add_text(s, "🗺️  Contexte", Inches(0.7), Inches(1.7), Inches(5.4), Inches(0.5),
         font_size=20, color=GREEN_DARK, bold=True)
ctx_items = [
    "Infrastructures urbaines dégradées (routes, éclairage)",
    "Problèmes d'assainissement et déchets sauvages",
    "Risques environnementaux (inondations, feux)",
    "Absence de plateforme centralisée de signalement",
    "Signalements informels (téléphone, réseaux sociaux)",
]
for i, t in enumerate(ctx_items):
    add_text(s, f"▸ {t}", Inches(0.9), Inches(2.4 + i * 0.55), Inches(5.2), Inches(0.5),
             font_size=13, color=DARK)
# Right column
add_card(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(5.0), fill=BLUE_LIGHT)
add_text(s, "❓  Problématique", Inches(7.0), Inches(1.7), Inches(5.6), Inches(0.5),
         font_size=20, color=BLUE, bold=True)
add_text(s, "Comment permettre aux citoyens mauritaniens de signaler efficacement les problèmes urbains et d'en assurer le suivi ?", Inches(7.0), Inches(2.4), Inches(5.6), Inches(1.0),
         font_size=15, color=DARK, bold=True)
objectives = [
    "Application mobile accessible (3 langues)",
    "Signalement en 4 étapes simples",
    "Carte interactive des signalements",
    "Système de crédibilité (votes)",
    "Messagerie intégrée temps réel",
    "Dashboard d'administration web",
]
for i, t in enumerate(objectives):
    add_text(s, f"✓ {t}", Inches(7.0), Inches(3.5 + i * 0.5), Inches(5.6), Inches(0.45),
             font_size=13, color=DARK)
add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — OBJECTIFS
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Objectifs du Projet", "Cahier des charges du module Développement Mobile")
goals = [
    ("👥", "Travail en groupe", "Participation active des 4 membres"),
    ("⚡", "Fonctionnalités interactives", "Formulaires, navigation, CRUD"),
    ("📦", "Gestion d'état (Provider)", "9 ChangeNotifier controllers"),
    ("🏗️", "Architecture MVC", "Séparation Modèle/Vue/Contrôleur"),
    ("🗄️", "Base de données + API", "Supabase PostgreSQL + REST + Realtime"),
    ("🌍", "Internationalisation", "3 langues (AR/FR/EN)"),
    ("✨", "Code structuré", "Organisation logique des fichiers"),
    ("🎯", "Projet original", "Adapté au contexte mauritanien"),
]
for i, (emoji, title, desc) in enumerate(goals):
    row = i // 4
    col = i % 4
    x = Inches(0.4 + col * 3.2)
    y = Inches(1.6 + row * 2.6)
    add_card(s, x, y, Inches(2.95), Inches(2.2))
    add_circle(s, Inches(x.inches + 0.85), Inches(y.inches + 0.2), Inches(0.9), fill_color=GREEN)
    add_text(s, emoji, Inches(x.inches + 0.85), Inches(y.inches + 0.25), Inches(0.9), Inches(0.8),
             font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(s, title, Inches(x.inches + 0.15), Inches(y.inches + 1.2), Inches(2.65), Inches(0.4),
             font_size=14, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(x.inches + 0.15), Inches(y.inches + 1.6), Inches(2.65), Inches(0.4),
             font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE MVC
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Architecture MVC", "Modèle — Vue — Contrôleur avec couche Services et DAOs")

# Diagram boxes
layers_data = [
    ("VUES\n(Flutter Widgets)", "13 écrans\nHomeView · MapView · ChatView\nAddReportView · SettingsView …", GREEN_LIGHT, GREEN),
    ("CONTROLEURS\n(ChangeNotifier)", "9 providers\nReport · Auth · Chat · Alert\nMap · AddReport · Theme · Locale · Nav", YELLOW_LIGHT, YELLOW),
    ("SERVICES / DAOs\n(Abstraction)", "Couche métier\nReportService · AuthService\nReportDao · UserDao · VoteDao …", BLUE_LIGHT, BLUE),
    ("SUPABASE\n(Cloud Backend)", "PostgreSQL + RLS\nAuth · Storage · Realtime\nFonctions SECURITY DEFINER", PURPLE_LIGHT, PURPLE),
]

for i, (title, desc, bg, border) in enumerate(layers_data):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.6)
    add_card(s, x, y, Inches(2.9), Inches(3.5), fill=bg)
    # Top colored bar
    add_rect_sharp(s, x, y, Inches(2.9), Inches(0.55), fill_color=border)
    add_text(s, title, x, y, Inches(2.9), Inches(0.55),
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, desc, Inches(x.inches + 0.15), Inches(y.inches + 0.7), Inches(2.6), Inches(2.6),
             font_size=11, color=DARK, alignment=PP_ALIGN.CENTER)

# Arrows between boxes
for i in range(3):
    x = Inches(3.3 + i * 3.2)
    y = Inches(3.0)
    add_text(s, "⟶", x, y, Inches(0.5), Inches(0.5), font_size=28, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Flow description
add_text(s, "← Flutter Provider : notifyListeners() → Selector / Consumer → Rebuild Widget →", Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
         font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# Tech stack badges
tech_badges = ["Flutter 3.3+", "Dart 3.x", "Provider 6.x", "Supabase 2.x", "PostgreSQL 15", "flutter_map"]
for i, badge in enumerate(tech_badges):
    x = Inches(0.4 + i * 2.15)
    add_badge(s, badge, x, Inches(6.0), Inches(1.95), Inches(0.6), color=GREEN, font_color=WHITE, font_size=12)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — STACK TECHNO
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Stack Technologique")

capture_path = os.path.join(DIAGRAMS, "stack-technologique.png")
if os.path.exists(capture_path):
    add_img(s, capture_path, Inches(1.5), Inches(1.6), width=Inches(10.3))
else:
    tech_stack = [
        ("📱", "Flutter / Dart", "Framework & Langage"),
        ("🎨", "Material 3", "Design System"),
        ("📦", "Provider", "Gestion d'état"),
        ("🗺️", "flutter_map + OSM", "Cartographie"),
        ("🗄️", "Supabase", "Backend (DB/Auth/Storage)"),
        ("📡", "Supabase Realtime", "Messages temps réel"),
        ("📷", "image_picker", "Photos signalement"),
        ("🌍", "intl + ARB", "Internationalisation (3 langues)"),
        ("⚙️", "geolocator", "Géolocalisation GPS"),
        ("🔗", "Vercel", "Déploiement Admin Web"),
    ]
    for i, (emoji, name, desc) in enumerate(tech_stack):
        col = i % 5
        row = i // 5
        x = Inches(0.4 + col * 2.55)
        y = Inches(1.7 + row * 2.6)
        add_card(s, x, y, Inches(2.3), Inches(2.2))
        add_text(s, emoji, Inches(x.inches + 0.3), Inches(y.inches + 0.15), Inches(1.7), Inches(0.7),
                 font_size=32, alignment=PP_ALIGN.CENTER)
        add_text(s, name, Inches(x.inches + 0.1), Inches(y.inches + 0.9), Inches(2.1), Inches(0.5),
                 font_size=13, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s, desc, Inches(x.inches + 0.1), Inches(y.inches + 1.4), Inches(2.1), Inches(0.5),
                 font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — BASE DE DONNÉES
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Base de Données", "Supabase PostgreSQL — 5 tables avec Row Level Security (RLS)")

db_tables = [
    ("users", "id, username, email, created_at, reputation", "Profils et authentification"),
    ("reports", "id, user_id, category, description, lat, lng, photo_url, status, confirm_count, deny_count", "Signalements citoyens"),
    ("votes", "id, report_id, user_id, vote_type (confirm/deny), created_at\nUNIQUE(report_id, user_id)", "Votes de crédibilité"),
    ("notifications", "id, user_id, report_id, message, is_read, created_at", "Notifications DB"),
    ("messages", "id, report_id, sender_id, receiver_id, content, is_read, created_at", "Messagerie temps réel"),
]

# Table header
col_widths = [Inches(1.5), Inches(6.0), Inches(4.0)]
col_starts = [Inches(0.5), Inches(2.0), Inches(8.0)]
headers = ["Table", "Colonnes principales", "Rôle"]
for i, (hdr, cs, cw) in enumerate(zip(headers, col_starts, col_widths)):
    add_rect_sharp(s, cs, Inches(1.6), cw, Inches(0.5), fill_color=GREEN)
    add_text(s, hdr, cs, Inches(1.6), cw, Inches(0.5),
             font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

for i, (table, cols, desc) in enumerate(db_tables):
    y = Inches(2.1 + i * 0.85)
    bg = GRAY_LIGHT if i % 2 == 0 else WHITE
    for j, (text, cs, cw) in enumerate(zip([table, cols, desc], col_starts, col_widths)):
        add_rect_sharp(s, cs, y, cw, Inches(0.85), fill_color=bg)
        fs = 12 if j == 0 else 11
        bd = j == 0
        add_text(s, text, Inches(cs.inches + 0.1), y, Inches(cw.inches - 0.2), Inches(0.85),
                 font_size=fs, color=DARK, bold=bd, anchor=MSO_ANCHOR.MIDDLE)

# RLS note
add_card(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8), fill=GREEN_LIGHT)
add_text(s, "🔒 Sécurité : RLS (Row Level Security) sur toutes les tables · Fonctions SECURITY DEFINER pour les votes et notifications · Authentification Supabase (JWT) · Bucket Storage avec politiques d'accès",
         Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.7), font_size=12, color=GREEN_DARK)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — FONCTIONNALITÉS
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Fonctionnalités Développées")

features_grid = [
    ("📝", "Signalement\n4 étapes", "Catégorie → Description →\nPhoto → Révision & envoi"),
    ("🗺️", "Carte Interactive\nOSM", "Marqueurs colorés par\ncatégorie · Cache FMTC · GPS"),
    ("👍", "Vote de\nCrédibilité", "Confirmer/Infirmer\nBarre de progression · UNIQUE"),
    ("💬", "Messagerie\nTemps Réel", "Par signalement · Realtime\nBadge non lus · Bulles"),
    ("🔔", "Notifications\nDB", "Créées à chaque signalement\nMarquer tout lu · Badge"),
    ("🌍", "3 Langues\nAR/FR/EN", "Internationalisation complète\n95+ clés · Persistance"),
    ("🌙", "Thème\nClair/Sombre", "Material 3 · Persistance\nMode système automatique"),
    ("📊", "Dashboard\nAdmin Web", "Vercel · Gestion signalements\nUtilisateurs · Chart.js"),
    ("📷", "Photos\nSignalement", "Caméra/Galerie · Supabase Storage\nURL publique dans report"),
]
for i, (emoji, title, desc) in enumerate(features_grid):
    row = i // 3
    col = i % 3
    x = Inches(0.4 + col * 4.2)
    y = Inches(1.5 + row * 1.9)
    add_card(s, x, y, Inches(3.9), Inches(1.65))
    add_circle(s, Inches(x.inches + 0.15), Inches(y.inches + 0.15), Inches(0.7), fill_color=GREEN)
    add_text(s, emoji, Inches(x.inches + 0.15), Inches(y.inches + 0.18), Inches(0.7), Inches(0.65),
             font_size=22, alignment=PP_ALIGN.CENTER)
    add_text(s, title, Inches(x.inches + 1.0), Inches(y.inches + 0.1), Inches(2.7), Inches(0.6),
             font_size=14, color=GREEN, bold=True)
    add_text(s, desc, Inches(x.inches + 1.0), Inches(y.inches + 0.7), Inches(2.7), Inches(0.8),
             font_size=10, color=GRAY)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — WIZARD 4 ÉTAPES
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Création d'un Signalement — Wizard 4 Étapes")

steps = [
    ("1", "Catégorie +\nDescription", "12 catégories\n280 caractères max\nValidation requise", GREEN_LIGHT, GREEN),
    ("2", "Localisation\nGéographique", "Carte OSM interactive\nPin central fixe\nGPS automatique", BLUE_LIGHT, BLUE),
    ("3", "Photo\n(Optionnelle)", "Caméra ou Galerie\nUpload Supabase Storage\nURL publique", YELLOW_LIGHT, YELLOW),
    ("4", "Révision\n& Envoi", "Récapitulatif complet\nbuildDraft() → addReport()\nSnackBar confirmation", PURPLE_LIGHT, PURPLE),
]

for i, (num, title, desc, bg, border) in enumerate(steps):
    x = Inches(0.3 + i * 3.25)
    y = Inches(1.5)
    add_card(s, x, y, Inches(2.95), Inches(4.2), fill=bg)
    # Step number
    add_circle(s, Inches(x.inches + 0.95), Inches(y.inches + 0.15), Inches(0.8), fill_color=border)
    add_text(s, num, Inches(x.inches + 0.95), Inches(y.inches + 0.15), Inches(0.8), Inches(0.8),
             font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title
    add_text(s, title, Inches(x.inches + 0.1), Inches(y.inches + 1.1), Inches(2.75), Inches(0.8),
             font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text(s, desc, Inches(x.inches + 0.2), Inches(y.inches + 2.0), Inches(2.55), Inches(1.5),
             font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    # Screenshot placeholder
    cap_name = [None, "nouveau-signalement.png", "carte-maps.png", "prendre-photo.png", "envoie-signalement.png"][i + 1]
    cap_path = os.path.join(CAPTURES, cap_name) if cap_name else None
    if cap_path and os.path.exists(cap_path):
        add_img(s, cap_path, Inches(x.inches + 0.35), Inches(y.inches + 3.1), width=Inches(2.0))

# Arrows between steps
for i in range(3):
    x = Inches(3.25 + i * 3.25)
    add_text(s, "⟶", x, Inches(3.3), Inches(0.5), Inches(0.5), font_size=24, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — CAPTURES D'ÉCRAN
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Aperçu de l'Application")

screenshots_map = [
    ("splash.png", "Splash"),
    ("connexion.png", "Connexion"),
    ("register.png", "Inscription"),
    ("acceuil.png", "Accueil"),
    ("carte-maps.png", "Carte"),
    ("nouveau-signalement.png", "Signalement"),
    ("prendre-photo.png", "Photo"),
    ("envoie-signalement.png", "Envoi"),
    ("detail_du_signalement.png", "Détail"),
    ("messagerie.png", "Messagerie"),
    ("alert.png", "Alertes"),
    ("profile.png", "Profil"),
    ("parametre.png", "Paramètres"),
    ("admin-dashboard.png", "Admin"),
]

for i, (fname, label) in enumerate(screenshots_map):
    cols = 4
    row = i // cols
    col = i % cols
    x = Inches(0.2 + col * 3.3)
    y = Inches(1.45 + row * 1.45)
    path = os.path.join(CAPTURES, fname)
    if os.path.exists(path):
        add_img(s, path, Inches(x.inches + 0.55), y, height=Inches(1.15))
        add_text(s, label, x, Inches(y.inches + 1.2), Inches(3.0), Inches(0.3),
                 font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    else:
        add_card(s, Inches(x.inches + 0.15), y, Inches(2.8), Inches(1.4), fill=GRAY_LIGHT)
        add_text(s, label, x, Inches(y.inches + 0.3), Inches(3.0), Inches(0.7),
                 font_size=16, color=GRAY_MED, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — NAVIGATION
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Navigation & Architecture de l'Application")

nav_path = os.path.join(DIAGRAMS, "navigation-app.png")
if os.path.exists(nav_path):
    add_img(s, nav_path, Inches(0.3), Inches(1.5), width=Inches(12.7))
else:
    # Fallback textual navigation
    add_card(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5), fill=GREEN_LIGHT)
    add_text(s, "🔄 Shell Principal (MainLayout)", Inches(0.7), Inches(1.7), Inches(5.4), Inches(0.4),
             font_size=16, color=GREEN_DARK, bold=True)
    tabs = ["Accueil (HomeView)", "Carte (MapView)", "Alertes (AlertsView)", "Compte (AccountView)"]
    for i, t in enumerate(tabs):
        add_text(s, f"  {i+1}. {t}", Inches(0.9), Inches(2.2 + i * 0.4), Inches(5.2), Inches(0.35),
                 font_size=13, color=DARK)

    add_card(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.5), fill=BLUE_LIGHT)
    add_text(s, "📌 Routes Push", Inches(7.0), Inches(1.7), Inches(5.4), Inches(0.4),
             font_size=16, color=BLUE, bold=True)
    routes = ["AddReportView (FAB +)", "ReportDetailView (depuis carte/liste)", "ChatView (depuis détail)", "EmergencyNumbersView (SOS/Compte)"]
    for i, r in enumerate(routes):
        add_text(s, f"  › {r}", Inches(7.2), Inches(2.2 + i * 0.4), Inches(5.2), Inches(0.35),
                 font_size=13, color=DARK)

    add_card(s, Inches(0.5), Inches(4.5), Inches(12.1), Inches(1.5), fill=YELLOW_LIGHT)
    add_text(s, "⚙️ Contrôle de Navigation", Inches(0.7), Inches(4.6), Inches(11.7), Inches(0.4),
             font_size=16, color=DARK, bold=True)
    nav_items = [
        "NavigationProvider (ChangeNotifier) → AppTab (enum 0-3)",
        "IndexedStack : tous les onglets vivants simultanément (pas de rebuild au switch)",
        "BottomAppBar avec notch pour le FAB central + animation pulsante",
        "Selector<NavigationProvider, int> : seul l'IndexedStack et la BottomNav se rebuildent",
    ]
    for i, item in enumerate(nav_items):
        add_text(s, f"  ▸ {item}", Inches(0.9), Inches(5.05 + i * 0.3), Inches(11.5), Inches(0.3),
                 font_size=11, color=DARK)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — FLUX SIGNALEMENT
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Flux de Données — Création d'un Signalement")

flow_path = os.path.join(DIAGRAMS, "flow-signalement.png")
if os.path.exists(flow_path):
    add_img(s, flow_path, Inches(0.3), Inches(1.5), width=Inches(12.7))
else:
    # Textual flow
    flow_steps = [
        ("👤", "Utilisateur", "Remplit le formulaire en 4 étapes"),
        ("→", "AddReportProvider", "Stocke l'état local (ChangeNotifier)"),
        ("→", "ReportProvider", "Appelle addReport() avec buildDraft()"),
        ("→", "ReportServiceDb", "DAO → INSERT dans Supabase"),
        ("→", "PostgreSQL", "RLS vérifie auth.uid() = user_id"),
        ("→", "NotificationService", "Crée notification pour tous les autres users"),
        ("→", "UI Mise à jour", "notifyListeners() → HomeView + MapView rebuild"),
    ]
    for i, (emoji, title, desc) in enumerate(flow_steps):
        y = Inches(1.6 + i * 0.75)
        add_card(s, Inches(0.5), y, Inches(12.3), Inches(0.6), fill=GRAY_LIGHT if i % 2 == 0 else WHITE)
        add_text(s, f"{emoji}  {title}", Inches(0.7), y, Inches(3.0), Inches(0.6),
                 font_size=13, color=DARK if i % 2 == 1 else GRAY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, Inches(4.0), y, Inches(8.5), Inches(0.6),
                 font_size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# Bottom details
add_card(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8), fill=GREEN_LIGHT)
add_text(s, "📦 AddReportProvider : currentStepIndex · selectedCategory · description · location · photoUrl · buildDraft()",
         Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.7), font_size=12, color=GREEN_DARK)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — DASHBOARD ADMIN
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Dashboard Administration (Web)")
# Left: features
add_card(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.0), fill=GREEN_LIGHT)
add_text(s, "🖥️  Fonctionnalités Admin", Inches(0.7), Inches(1.7), Inches(5.4), Inches(0.5),
         font_size=18, color=GREEN_DARK, bold=True)
admin_features = [
    "Login via Supabase Auth",
    "Gestion des signalements (filtre, statut, suppression)",
    "Gestion des utilisateurs (stats, rôle admin)",
    "Graphique signalements par catégorie (Chart.js)",
    "Notifications avec marquage lecture",
    "Internationalisation AR/FR",
    "Support RTL/LTR",
]
for i, feat in enumerate(admin_features):
    add_text(s, f"✓  {feat}", Inches(0.9), Inches(2.4 + i * 0.55), Inches(5.2), Inches(0.5),
             font_size=13, color=DARK)
# Right: screenshot
admin_cap = os.path.join(CAPTURES, "admin-dashboard.png")
if os.path.exists(admin_cap):
    add_img(s, admin_cap, Inches(6.6), Inches(1.5), width=Inches(6.2))
else:
    add_card(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.0), fill=PURPLE_LIGHT)
    add_text(s, "Dashboard Admin", Inches(7.0), Inches(3.5), Inches(5.4), Inches(0.5),
             font_size=18, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
# URL
add_card(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.7), fill=YELLOW_LIGHT)
add_text(s, "🔗  admin-dashboard-pearl-delta-63.vercel.app", Inches(1.0), Inches(5.55), Inches(11.1), Inches(0.6),
         font_size=15, color=DARK, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 15 — SÉCURITÉ
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Sécurité", "Protection des données à plusieurs niveaux")

sec_items = [
    ("🔒", "Row Level Security (RLS)", "Politiques par table : chaque utilisateur ne voit/modifie que ses propres données", GREEN),
    ("🛡️", "SECURITY DEFINER", "Fonctions RPC contournant RLS pour les opérations critiques (votes, notifications)", BLUE),
    ("🔑", "Authentification Supabase", "JWT tokens · Sessions persistantes · Email + Google OAuth 2.0", PURPLE),
    ("📁", "Stockage Sécurisé", "Bucket 'reports' avec politiques : upload = auth, lecture = tout auth, supprimer = owner", YELLOW),
    ("🔐", "Validation Côté Client", "Contrôle des formulaires · 280 caractères max · Catégorie requise · Vote unique (UNIQUE contrainte)", GREEN),
]

for i, (emoji, title, desc, accent) in enumerate(sec_items):
    y = Inches(1.6 + i * 1.05)
    add_card(s, Inches(0.5), y, Inches(12.3), Inches(0.9), fill=GRAY_LIGHT)
    add_circle(s, Inches(0.65), Inches(y.inches + 0.1), Inches(0.65), fill_color=accent)
    add_text(s, emoji, Inches(0.65), Inches(y.inches + 0.12), Inches(0.65), Inches(0.6),
             font_size=22, alignment=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.5), Inches(y.inches + 0.05), Inches(4.0), Inches(0.4),
             font_size=15, color=DARK, bold=True)
    add_text(s, desc, Inches(1.5), Inches(y.inches + 0.45), Inches(11.0), Inches(0.4),
             font_size=12, color=GRAY)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 16 — RESPECT CONSIGNES
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, WHITE)
add_header_bar(s, "Respect du Cahier des Charges")

checks = [
    "👥 Travail en groupe (4 membres)",
    "⚡ Fonctionnalités interactives (formulaires, navigation, CRUD)",
    "📊 Gestion d'état (Provider + 9 ChangeNotifier)",
    "🏗️ Architecture MVC (Modèles / Vues / Contrôleurs + Services)",
    "🗄️ Base de données (Supabase PostgreSQL)",
    "🌐 API (Supabase REST + Realtime + Storage)",
    "🌍 Internationalisation (3 langues : AR/FR/EN)",
    "✅ Code structuré et maintenable",
    "🎯 Projet original adapté au contexte mauritanien",
]
for i, check in enumerate(checks):
    row = i // 3
    col = i % 3
    x = Inches(0.4 + col * 4.2)
    y = Inches(1.6 + row * 1.7)
    card = add_card(s, x, y, Inches(3.9), Inches(1.4), fill=GREEN_LIGHT)
    add_text(s, check, Inches(x.inches + 0.2), Inches(y.inches + 0.2), Inches(3.5), Inches(1.0),
             font_size=14, color=GREEN_DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 17 — MERCI
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, GREEN)

add_text(s, "Merci de votre attention", Inches(1.5), Inches(1.0), Inches(10.3), Inches(1.2),
         font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect_sharp(s, Inches(5.5), Inches(2.1), Inches(2.3), Inches(0.05), fill_color=YELLOW)

# Repository
add_text(s, "📂  Code source", Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.5),
         font_size=18, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "github.com/mohameden19961/project-baligh", Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.5),
         font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

# Dashboard
add_text(s, "📊  Dashboard Admin", Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.5),
         font_size=18, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "admin-dashboard-pearl-delta-63.vercel.app", Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
         font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

# Separator
add_rect_sharp(s, Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.05), fill_color=YELLOW)

# Questions
add_text(s, "Des questions ?", Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.8),
         font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "Institut Supérieur du Numérique (SUPNUM) — L2 DWM — 2025-2026",
         Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
         font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(s)

# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
output_path = "baligh_presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path} ({slide_num[0]} slides)")
